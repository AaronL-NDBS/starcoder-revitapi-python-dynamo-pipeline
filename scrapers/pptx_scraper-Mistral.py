# pptx_scraper.py
import json, os
from pathlib import Path
from pptx import Presentation

RAW_DIR = Path("dataset/pptx/raw")
RAW_DIR.mkdir(parents=True, exist_ok=True)

def scrape_pptx_folder(folder_path):
    files = list(Path(folder_path).rglob("*.pptx"))
    for p in files:
        print(f"Scraping: {p.name}")
        try:
            prs = Presentation(p)
            slides_data = []
            for i, slide in enumerate(prs.slides):
                text = "\n".join([s.text for s in slide.shapes if hasattr(s, "text")])
                notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
                slides_data.append({"slide": i+1, "content": text, "notes": notes})
            
            with open(RAW_DIR / f"{p.stem}.json", "w") as f:
                json.dump({"filename": p.name, "slides": slides_data}, f)
        except Exception as e:
            print(f"  Error: {e}")

if __name__ == "__main__":
    path = input("Enter path to PPTX folder: ").strip('"')
    scrape_pptx_folder(path)