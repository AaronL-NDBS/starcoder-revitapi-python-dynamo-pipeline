# tools/pptx_scraper.py
import json
import os
import re
import shutil
import subprocess
import time
import zipfile
import requests
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    exit(1)

# Paths
RAW_DIR      = Path("dataset/pptx/raw")
APPROVED_DIR = Path("dataset/pptx/approved")
CLONE_DIR    = Path("dataset/pptx/cloned_repos")
PROGRESS_LOG = Path("dataset/pptx/progress_log.json")

for d in [RAW_DIR, APPROVED_DIR, CLONE_DIR]:
    d.mkdir(parents=True, exist_ok=True)

OLLAMA_URL   = "http://localhost:11434/api/chat"
JUDGE_MODEL  = "mistral"
MIN_SCORE    = 5

# Single-Pass System Prompt
SINGLE_PASS_SYSTEM = """You extract training data for a Revit API and Dynamo Python code model.
Analyze the PowerPoint slide content. If it contains technical Revit/Dynamo concepts, score its relevance (0-10) and extract 1-2 training pairs.
Pairs must use the Dynamo IN[]/OUT pattern and include clr.AddReference('RevitAPI').

Return ONLY valid JSON in this exact format:
{
  "score": <0-10>,
  "reason": "<brief justification>",
  "pairs": [
    {"prompt": "<clear question/task>", "completion": "<python code>"}
  ]
}
If no extractable code exists, return {"score": 0, "reason": "No technical content", "pairs": []}"""

def load_progress():
    if PROGRESS_LOG.exists():
        with open(PROGRESS_LOG, "r") as f:
            return set(json.load(f))
    return set()

def save_progress(processed_files):
    with open(PROGRESS_LOG, "w") as f:
        json.dump(list(processed_files), f, indent=4)

def extract_pptx(pptx_path: Path) -> list[dict]:
    slides = []
    try:
        prs = Presentation(str(pptx_path))
        for i, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip(): text_parts.append(para.text.strip())

            notes = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide and slide.notes_slide.notes_text_frame else ""
            slide_text = "\n".join(text_parts)
            
            if slide_text or notes:
                slides.append({
                    "slide_number": i + 1,
                    "combined": f"{slide_text}\n\nSpeaker notes:\n{notes}" if notes else slide_text
                })
    except Exception as e:
        print(f"    Error reading {pptx_path.name}: {e}")
    return slides

def process_slide_single_pass(slide_content: str, source_name: str) -> dict:
    try:
        r = requests.post(OLLAMA_URL, json={
            "model": JUDGE_MODEL,
            "messages": [
                {"role": "system", "content": SINGLE_PASS_SYSTEM},
                {"role": "user", "content": f"SOURCE: {source_name}\n\nSLIDE CONTENT:\n{slide_content[:1500]}"}
            ],
            "stream": False
        }, timeout=120)
        content = r.json()["message"]["content"]
        match = re.search(r'\{.*\}', content, re.DOTALL)
        if match:
            return json.loads(match.group())
    except Exception as e:
        print(f"    Mistral error: {e}")
    return {"score": 0, "pairs": []}

def process_pptx(pptx_path: Path, processed_files: set) -> dict:
    safe_name = re.sub(r'[^\w]', '_', pptx_path.stem)[:60]
    
    if safe_name in processed_files:
        print(f"  [SKIPPED] {pptx_path.name} (Already processed)")
        return {"file": pptx_path.name, "approved": 0}

    print(f"\n  [PROCESSING] {pptx_path.name}")
    slides = extract_pptx(pptx_path)
    
    if not slides:
        processed_files.add(safe_name)
        save_progress(processed_files)
        return {"file": pptx_path.name, "approved": 0}

    out_path = APPROVED_DIR / f"{safe_name}.jsonl"
    approved_pairs = []

    for idx, slide in enumerate(slides):
        print(f"    -> Slide {idx+1}/{len(slides)}", end="\r")
        result = process_slide_single_pass(slide["combined"], pptx_path.name)
        
        if result.get("score", 0) >= MIN_SCORE:
            for pair in result.get("pairs", []):
                if pair.get("prompt") and pair.get("completion"):
                    record = {
                        "source": "pptx",
                        "file": pptx_path.name,
                        "slide": slide["slide_number"],
                        "prompt": pair["prompt"],
                        "completion": pair["completion"],
                        "date_added": datetime.now().strftime("%Y-%m-%d")
                    }
                    approved_pairs.append(record)
                    with open(out_path, "a", encoding="utf-8") as f:
                        f.write(json.dumps(record) + "\n")

    processed_files.add(safe_name)
    save_progress(processed_files)
    print(f"    [DONE] Extracted {len(approved_pairs)} valid pairs.    ")
    return {"file": pptx_path.name, "approved": len(approved_pairs)}

def main():
    print("=" * 60)
    print(f" PowerPoint Scraper | Model: {JUDGE_MODEL}")
    print("=" * 60)
    
    processed_files = load_progress()
    raw = input("\nPath to .pptx file or folder (or 'q' to quit): ").strip().strip('"')
    if raw.lower() == 'q': return
    
    path = Path(raw)
    files = [path] if path.is_file() and path.suffix.lower() == ".pptx" else sorted(path.rglob("*.pptx"))
    
    to_process = [f for f in files if re.sub(r'[^\w]', '_', f.stem)[:60] not in processed_files]
    print(f"\nFound {len(files)} files. ({len(files) - len(to_process)} already processed)")

    total_approved = 0
    for pptx_path in to_process:
        res = process_pptx(pptx_path, processed_files)
        total_approved += res.get("approved", 0)

    print("\n" + "=" * 60)
    print(f" Session Complete | High-confidence pairs: {total_approved}")
    print("=" * 60)

if __name__ == "__main__":
    main()